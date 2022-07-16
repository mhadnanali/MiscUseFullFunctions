import pandas as pd
from pptx import Presentation

def PresentationPPT():


        allwordsDF=pd.read_csv('HSK3PPT.csv',index_col=0)
        root = Presentation()
        # Creating slide layout
        f1 = root.slide_layouts[0] #Title Slide
        f2 = root.slide_layouts[2] #Section header
        f3 = root.slide_layouts[1] #title and content        
        f4 = root.slide_layouts[6]# blank for Picture
        """ Ref for slide types:
        0 -> title and subtitle
        1 -> title and content
        2 -> section header
        3 -> two content
        4 -> Comparison
        5 -> Title only
        6 -> Blank
        7 -> Content with caption
        8 -> Pic with caption
        """
        for index, row in allwordsDF.iterrows():
            #print(row['Word'], row['Pinyin'], row['Meaning'],row['Sentence'])                    
            slide = root.slides.add_slide(f2)
            slide.shapes.title.text = str(row['Word'])#+'\n'+row['Meaning']
            slide.placeholders[1].text =row['Pinyin']+'\n'+row['Meaning']+'\n('+str(row['POF'])+')'

        """
        slide4 = root.slides.add_slide(f3)
        slide4.shapes.title.text = "How Many? "
        slide4.placeholders[1].text = "Tweets:\t\t"+ str(self.allTweets)+"\nHashtags:\t"+str(self.allHash)+"\nUsers:\t\t\t"+str(self.allAuthors)+" \n\t\t" \
         "Verified Users:\t\t\t"+str(self.verifiedcounter)+"\n\t\tNon-Verified Users:\t"+str(self.allAuthors- self.verifiedcounter)


        #####
        left = top = Inches(0)
        slide15 = root.slides.add_slide(f4)
        slide15.shapes.add_picture(oPath+'1 Tweets of '+self.topHashtag+' Per Day.jpg', left,top)#, height = height)
        slide16 = root.slides.add_slide(f4)
        slide16.shapes.add_picture(oPath+'3 Users of '+self.topHashtag+' Per Day.jpg', left,top)#, height = height)
        slide9 = root.slides.add_slide(f4)
        slide9.shapes.add_picture(path+' 3 Word Count.jpg', left,top)#, height = height)
        slide8 = root.slides.add_slide(f4)
        slide8.shapes.add_picture(path+' 4 Top25.jpg', left,top)#, height = height)
        slide18 = root.slides.add_slide(f2)
        slide18.shapes.title.text = "If I get 500 Subscribers"
        slide18.placeholders[1].text = "\nTop 5 influential users of the Hashtag\nTop 5 Verified Accounts of " \
                                       "Hashtag\nLocation based Analysis of Hashtag\n5 Lowest intention getting Users"
        slide19 = root.slides.add_slide(f2)
        slide19.shapes.title.text = "If I get 1000 Subscribers"
        slide19.placeholders[1].text = "\nPerson based Analysis (e.g Imran Khan related Analysis)\nAnalysis of Twitter profiles"
        slide20 = root.slides.add_slide(f2)
        slide20.shapes.title.text = "Data Mining is time taking task... \n\t\tSo... subscribe and like :/"
          # Saving file
        """
        root.save("HSK3.pptx")

        print("done")

PresentationPPT()
